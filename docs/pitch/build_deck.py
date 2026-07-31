# -*- coding: utf-8 -*-
"""VLearn Interaction Layer - pitch deck generator (style: VLearn HTML deck)."""
import os
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r"E:\Labs\Hackathon-Team-DO-MESSIU\docs\pitch\VLearn-Pitch-Deck.pptx"

# ---------------------------------------------------------------- design tokens
FONT = "Segoe UI"
PRIMARY   = RGBColor(0x13, 0x4D, 0x8B)
PRIMARY_D = RGBColor(0x0B, 0x2F, 0x55)
PRIMARY_L = RGBColor(0x2B, 0x88, 0xFF)
ACCENT    = RGBColor(0xC7, 0x21, 0x27)
ACCENT_D  = RGBColor(0x8C, 0x11, 0x17)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_D  = RGBColor(0xB3, 0x6A, 0x07)
TEXT      = RGBColor(0x06, 0x1B, 0x49)
BODY      = RGBColor(0x30, 0x43, 0x5E)
MUTED     = RGBColor(0x52, 0x65, 0x7F)
SMALLC    = RGBColor(0x50, 0x66, 0x80)
LINE      = RGBColor(0xDD, 0xE8, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ONBLUE    = RGBColor(0xE2, 0xEF, 0xFF)
ONBLUE_D  = RGBColor(0x9E, 0xCB, 0xFF)
ONRED     = RGBColor(0xFF, 0xDA, 0xDB)
SHADOW    = RGBColor(0xD3, 0xE1, 0xF3)
SHADOW_B  = RGBColor(0x07, 0x2F, 0x5C)
BARBG     = RGBColor(0xE8, 0xF0, 0xFA)
FOOT      = RGBColor(0x7F, 0x96, 0xB4)
SLIDENO   = RGBColor(0x93, 0xA9, 0xC5)

W, H = 13.333, 7.5
M = 0.82
CW = W - 2 * M
TOP = 0.46
BODY_Y = 1.30
FOOT_Y = 6.88
TOTAL = 13

WARN = []
_CUR = [0]

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- text metrics
_FCACHE = {}
_REG = r"C:\Windows\Fonts\segoeui.ttf"
_BLD = r"C:\Windows\Fonts\segoeuib.ttf"
_PREC = 4  # render font at 4x pt for sub-pixel precision


def _pil(size, bold):
    key = (round(size, 1), bool(bold))
    if key not in _FCACHE:
        _FCACHE[key] = ImageFont.truetype(_BLD if bold else _REG, int(round(size * _PREC)))
    return _FCACHE[key]


def tw(text, size, bold=False):
    """Text width in inches."""
    return _pil(size, bold).getlength(text) / _PREC / 72.0


def _tokens(parts):
    out = []
    for t, b in parts:
        for i, seg in enumerate(t.split("\n")):
            if i:
                out.append(("\n", b))
            for j, wd in enumerate(seg.split(" ")):
                out.append((wd, b))
    return out


def wrap_lines(parts, width, size):
    """Number of rendered lines for a paragraph made of (text, bold) parts."""
    lines, cur = 1, 0.0
    space = None
    for wd, b in _tokens(parts):
        if wd == "\n":
            lines += 1
            cur = 0.0
            space = None
            continue
        ww = tw(wd, size, b)
        sp = tw(" ", size, b) if cur > 0 else 0.0
        if cur > 0 and cur + sp + ww > width + 0.008:
            lines += 1
            cur = ww
        else:
            cur += sp + ww
    return lines


# ---------------------------------------------------------------- primitives
def _noline(shape):
    shape.line.fill.background()
    return shape


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def _grad(shape, c1, c2, angle=90):
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    try:
        shape.fill.gradient_angle = angle
    except Exception:
        pass
    return shape


def _noshadow(shape):
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    return shape


def rect(slide, x, y, w, h, radius=0.11, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    s.shadow.inherit = False
    _noline(s)
    _noshadow(s)
    if s.has_text_frame:
        s.text_frame.word_wrap = True
    return s


def oval(slide, x, y, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    _noline(s)
    _noshadow(s)
    return s


def center_text(shape, text, size, color=WHITE, bold=True, line=1.15):
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


class Box(object):
    """Text box that measures itself and reports overflow."""

    def __init__(self, slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        self.x, self.y = x, y
        self.shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = self.shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        self.tf = tf
        self.w, self.h = w, h
        self.n = 0
        self.need = 0.0
        self.chrome = False

    def add(self, content, size=13.5, bold=False, color=BODY, align=PP_ALIGN.LEFT,
            line=1.2, before=0, after=0, space=0):
        para = self.tf.paragraphs[0] if self.n == 0 else self.tf.add_paragraph()
        self.n += 1
        para.alignment = align
        para.line_spacing = line
        para.space_before = Pt(before)
        para.space_after = Pt(after)
        parts = content if isinstance(content, (list, tuple)) else [(content, color, bold)]
        norm = []
        for part in parts:
            t = part[0]
            c = part[1] if len(part) > 1 and part[1] is not None else color
            b = part[2] if len(part) > 2 and part[2] is not None else bold
            norm.append((t, b))
            r = para.add_run()
            r.text = t
            f = r.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = b
            f.color.rgb = c
            if space:
                r.font._rPr.set("spc", str(int(space * 100)))
        eff_w = self.w if not space else self.w - sum(len(t) for t, _ in norm) * space / 72.0
        nl = wrap_lines(norm, max(eff_w, 0.4), size)
        self.need += nl * 1.2 * size * line / 72.0 + (before + after) / 72.0
        flat = "".join(t for t, _ in norm)[:52]
        if self.need > self.h + 0.03:
            WARN.append("S%02d overflow %.2f>%.2f  %r" % (_CUR[0], self.need, self.h, flat))
        if self.y + self.need > FOOT_Y - 0.14 and not self.chrome:
            WARN.append("S%02d below-fold y=%.2f end=%.2f  %r"
                        % (_CUR[0], self.y, self.y + self.need, flat))
        return self


def label(slide, x, y, w, h, *args, **kwargs):
    anchor = kwargs.pop("anchor", MSO_ANCHOR.TOP)
    chrome = kwargs.pop("chrome", False)
    b = Box(slide, x, y, w, h, anchor)
    b.chrome = chrome
    if args:
        b.add(*args, **kwargs)
    return b


def flow(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """Stacked paragraphs in one box: [(text,size,bold,color,line,before)]"""
    b = Box(slide, x, y, w, h, anchor)
    for p in paras:
        text, size, bold, color = p[0], p[1], p[2], p[3]
        line = p[4] if len(p) > 4 else 1.25
        before = p[5] if len(p) > 5 else 0
        align = p[6] if len(p) > 6 else PP_ALIGN.LEFT
        space = p[7] if len(p) > 7 else 0
        b.add(text, size, bold, color, align, line, before, space=space)
    return b


# ---------------------------------------------------------------- components
def card(slide, x, y, w, h, kind="white", lift=0.085, radius=0.11):
    if kind == "white":
        _solid(rect(slide, x, y + lift, w, h, radius), SHADOW)
        c = rect(slide, x, y, w, h, radius)
        _solid(c, WHITE)
        c.line.color.rgb = LINE
        c.line.width = Pt(1)
    elif kind == "ice":
        _solid(rect(slide, x, y + lift, w, h, radius), SHADOW)
        c = rect(slide, x, y, w, h, radius)
        _grad(c, WHITE, RGBColor(0xE7, 0xF1, 0xFF), 45)
        c.line.color.rgb = LINE
        c.line.width = Pt(1)
    elif kind == "blue":
        _solid(rect(slide, x, y + lift, w, h, radius), SHADOW_B)
        c = rect(slide, x, y, w, h, radius)
        _grad(c, RGBColor(0x0F, 0x64, 0xC7), RGBColor(0x09, 0x2C, 0x58), 45)
    elif kind == "red":
        _solid(rect(slide, x, y + lift, w, h, radius), ACCENT_D)
        c = rect(slide, x, y, w, h, radius)
        _grad(c, RGBColor(0xFF, 0x39, 0x3F), RGBColor(0x84, 0x10, 0x16), 45)
    elif kind == "orange":
        _solid(rect(slide, x, y + lift, w, h, radius), ORANGE_D)
        c = rect(slide, x, y, w, h, radius)
        _grad(c, RGBColor(0xFF, 0xB1, 0x3B), RGBColor(0xC2, 0x6F, 0x03), 45)
    return c


def pill(slide, x, y, text, size=12, kind="light", h=0.36, pad=0.26):
    w = tw(text, size, True) + pad * 2
    s = rect(slide, x, y, w, h, 0.5)
    if kind == "light":
        _solid(s, WHITE)
        s.line.color.rgb = RGBColor(0xDC, 0xE9, 0xF8)
        s.line.width = Pt(1)
        col = PRIMARY
    elif kind == "ghost":
        _solid(s, RGBColor(0xED, 0xF5, 0xFF))
        s.line.color.rgb = RGBColor(0xCD, 0xE0, 0xF6)
        s.line.width = Pt(1)
        col = PRIMARY
    elif kind == "red":
        _grad(s, RGBColor(0xFF, 0x31, 0x38), ACCENT, 90)
        col = WHITE
    else:
        _grad(s, RGBColor(0x23, 0x86, 0xFF), PRIMARY, 90)
        col = WHITE
    center_text(s, text, size, col)
    return s, w


def btn(slide, x, y, text, size=13, h=0.48, pad=0.32):
    w = tw(text, size, True) + pad * 2
    sh = rect(slide, x, y + 0.07, w, h, 0.28)
    _solid(sh, RGBColor(0x08, 0x37, 0x6C))
    b = rect(slide, x, y, w, h, 0.28)
    _grad(b, RGBColor(0x23, 0x86, 0xFF), PRIMARY, 90)
    center_text(b, text, size, WHITE)
    return b, w, sh


def icon(slide, x, y, glyph, size=0.44, color=PRIMARY, fsize=None):
    _solid(rect(slide, x, y + 0.05, size, size, 0.3), RGBColor(0xD8, 0xE7, 0xF8))
    i = rect(slide, x, y, size, size, 0.3)
    _grad(i, WHITE, RGBColor(0xED, 0xF6, 0xFF), 45)
    if fsize is None:
        fsize = 15 if len(glyph) <= 1 else 10
    center_text(i, glyph, fsize, color)
    i.text_frame.word_wrap = False
    return i


def numchip(slide, x, y, n, size=0.42, fsize=14):
    c = rect(slide, x, y, size, size, 0.3)
    _grad(c, RGBColor(0x25, 0x88, 0xFF), PRIMARY, 45)
    center_text(c, n, fsize, WHITE)
    return c


def bar(slide, x, y, w, pct, color="blue", h=0.16):
    _solid(rect(slide, x, y, w, h, 0.5), BARBG)
    fl = rect(slide, x, y, max(w * pct, 0.18), h, 0.5)
    if color == "blue":
        _grad(fl, RGBColor(0x0F, 0x7D, 0xF0), PRIMARY, 0)
    elif color == "red":
        _grad(fl, RGBColor(0xFF, 0x4B, 0x51), ACCENT, 0)
    elif color == "orange":
        _grad(fl, RGBColor(0xFF, 0xBC, 0x54), ORANGE, 0)
    else:
        _solid(fl, RGBColor(0xC3, 0xD3, 0xE6))
    return fl


def arrow(slide, x, y, w, color=PRIMARY, h=0.18):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    _noline(s)
    _noshadow(s)
    _solid(s, color)
    try:
        s.adjustments[0] = 0.42
        s.adjustments[1] = 0.55
    except Exception:
        pass
    return s


def logo(slide, x, y, size):
    scale = Emu(int(Inches(size) / 64))
    paths = [([(8, 10), (26, 26), (18, 53), (2, 28)], ACCENT),
             ([(26, 26), (57, 5), (57, 33), (19, 59)], PRIMARY),
             ([(26, 26), (57, 5), (43, 31), (19, 59)], PRIMARY_D)]
    for pts, color in paths:
        b = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale)
        b.add_line_segments(pts[1:], close=True)
        sh = b.convert_to_shape(Inches(x), Inches(y))
        sh.shadow.inherit = False
        _noline(sh)
        _noshadow(sh)
        _solid(sh, color)


def dots(slide, x, y, cols, rows, gap=0.19, d=0.05, color=RGBColor(0xBD, 0xD3, 0xF0)):
    for r in range(rows):
        for c in range(cols):
            _solid(oval(slide, x + c * gap, y + r * gap, d, d), color)


# ---------------------------------------------------------------- slide chrome
def new_slide(deco=True, wave=False):
    _CUR[0] += 1
    s = prs.slides.add_slide(BLANK)
    _grad(rect(s, -0.3, -0.3, W + 0.6, H + 0.6, 0, MSO_SHAPE.RECTANGLE),
          WHITE, RGBColor(0xE9, 0xF2, 0xFE), 45)
    if deco:
        _grad(oval(s, W - 2.7, -2.0, 4.7, 4.7), RGBColor(0xDF, 0xEE, 0xFF),
              RGBColor(0xF3, 0xF9, 0xFF), 45)
        _grad(oval(s, -2.4, H - 2.4, 4.7, 4.7), RGBColor(0xE1, 0xEE, 0xFE),
              RGBColor(0xF5, 0xFA, 0xFF), 45)
    if wave:
        _grad(oval(s, -1.5, H - 1.6, W + 3, 3.6), RGBColor(0xE6, 0xF2, 0xFF),
              RGBColor(0xF7, 0xFB, 0xFF), 90)
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def topline(slide, kicker, no):
    label(slide, M, TOP, 8.0, 0.32, kicker.upper(), 12, True, PRIMARY, space=0.9)
    label(slide, W - M - 3.0, TOP, 3.0, 0.32, "%02d" % no, 12, True, SLIDENO, align=PP_ALIGN.RIGHT)


def footer(slide, right_text):
    logo(slide, M, FOOT_Y, 0.2)
    label(slide, M + 0.3, FOOT_Y - 0.01, 3.0, 0.28, "VLearn", 11, True, PRIMARY, chrome=True)
    label(slide, W - M - 4.5, FOOT_Y - 0.01, 4.5, 0.28, right_text, 11, True, FOOT,
          align=PP_ALIGN.RIGHT, chrome=True)


def h2(slide, x, y, w, text, size=29, h=1.2):
    return label(slide, x, y, w, h, text, size, True, TEXT, line=1.08)


def sub(slide, x, y, w, text, size=13.5, h=0.9, color=MUTED):
    return label(slide, x, y, w, h, text, size, True, color, line=1.42)


def feature(slide, x, y, w, h, glyph, title, desc, tsize=14, dsize=11.5):
    card(slide, x, y, w, h)
    icon(slide, x + 0.28, y + 0.26, glyph)
    flow(slide, x + 0.28, y + 0.86, w - 0.56, h - 1.06,
         [(title, tsize, True, PRIMARY_D, 1.15, 0),
          (desc, dsize, False, MUTED, 1.32, 7)])


# ================================================================ 01 title
def slide1():
    s = new_slide()
    dots(s, W - 2.6, 0.6, 8, 5)
    logo(s, M, 0.44, 0.5)
    label(s, M + 0.66, 0.5, 4, 0.42,
          [("V", ACCENT, True), ("Learn", RGBColor(0x0F, 0x3F, 0x73), True)], 22)
    pl, pw = pill(s, 0, 0.55, "Lớp tương tác tích hợp LMS", 12)
    pl.left = Inches(W - M - pw)

    label(s, 0, 2.05, W, 0.3, "PITCH DECK  •  HACKATHON", 12, True, PRIMARY,
          align=PP_ALIGN.CENTER, space=0.9)
    label(s, 0, 2.45, W, 1.2, "VLearn Interaction Layer", 54, True, TEXT,
          align=PP_ALIGN.CENTER, line=1.02)
    label(s, W / 2 - 4.8, 3.72, 9.6, 1.0,
          "Sinh viên không cần thêm một nơi để đọc tài liệu.\nHọ cần một cách để được lắng nghe ngay trong lúc học.",
          16, True, MUTED, align=PP_ALIGN.CENTER, line=1.5)

    b, bw, bsh = btn(s, 0, 4.95, "Xem lớp đang hiểu đến đâu", 14)
    p2, pw2 = pill(s, 0, 5.02, "Human-in-the-loop", 13, "red", h=0.41)
    total = bw + 0.22 + pw2
    b.left = bsh.left = Inches(W / 2 - total / 2)
    p2.left = Inches(W / 2 - total / 2 + bw + 0.22)
    footer(s, "01 / %d" % TOTAL)


# ================================================================ 02 problem
def slide2():
    s = new_slide()
    topline(s, "01. Vấn đề", 2)
    h2(s, M, BODY_Y, 6.1, "Cả lớp im lặng — nhưng\nim lặng không có nghĩa\nlà đã hiểu.", 30, 1.65)
    sub(s, M, BODY_Y + 1.78, 5.95,
        "Giảng viên trình chiếu PDF và tiếp tục giảng. Phía dưới, có người đã hiểu, có người chỉ hiểu một phần, "
        "và có người đã kẹt từ vài slide trước. Gần như không ai nói ra.", 13.5, 1.0)

    ay, aw = 4.28, 5.95
    cw = (aw - 0.4) / 3
    items = [("?", "Ngại phát biểu", "Sợ làm gián đoạn tiết học."),
             ("!", "Sợ câu hỏi đơn giản", "Không muốn bị đánh giá."),
             ("~", "Không biết mô tả", "Chưa gọi tên được điều mình chưa hiểu.")]
    for i, (g, t, d) in enumerate(items):
        x = M + i * (cw + 0.2)
        card(s, x, ay, cw, 2.15)
        icon(s, x + 0.24, ay + 0.24, g, 0.4, fsize=14)
        flow(s, x + 0.24, ay + 0.78, cw - 0.48, 1.32,
             [(t, 12.5, True, PRIMARY_D, 1.15, 0), (d, 11, False, MUTED, 1.3, 6)])

    bx = 7.5
    bw2 = W - M - bx
    card(s, bx, BODY_Y - 0.02, bw2, 4.98, "ice")
    label(s, bx + 0.44, BODY_Y + 0.34, bw2 - 0.88, 1.7,
          "“Giảng viên nhìn cả lớp im lặng và mặc định rằng mọi người vẫn đang theo kịp.”",
          19, True, TEXT, line=1.3)
    label(s, bx + 0.44, BODY_Y + 2.22, bw2 - 0.88, 0.28, "MỨC ĐỘ HIỂU THẬT CỦA LỚP", 10,
          True, SMALLC, space=0.8)
    rows = [("Đã hiểu", 0.58, "blue"), ("Hơi mơ hồ", 0.27, "orange"), ("Đang bị kẹt", 0.15, "red")]
    y = BODY_Y + 2.62
    for name, pct, col in rows:
        label(s, bx + 0.44, y - 0.02, 1.5, 0.26, name, 11.5, True, PRIMARY_D)
        bar(s, bx + 1.9, y + 0.02, bw2 - 2.75, pct, col, 0.15)
        y += 0.46
    label(s, bx + 0.44, y + 0.12, bw2 - 0.88, 0.8,
          [("Nhưng trên màn hình của giảng viên: ", MUTED, False),
           ("không có dữ liệu nào.", ACCENT, True)], 12.5, line=1.35)
    footer(s, "Problem")


# ================================================================ 03 evidence
def slide3():
    s = new_slide()
    topline(s, "02. Bằng chứng", 3)
    h2(s, M, BODY_Y, 5.5, "Khảo sát 36 sinh viên:\npain là có thật.", 30, 1.15)
    sub(s, M, BODY_Y + 1.28, 5.3,
        "Sinh viên có nhu cầu được hỗ trợ. Vấn đề là kênh hỗ trợ hiện tại không dùng được ngay trong giờ học.",
        13.5, 0.9)

    py = 3.5
    card(s, M, py, 5.3, 1.55, "blue")
    label(s, M + 0.44, py + 0.26, 4.4, 0.9, "88,9%", 42, True, WHITE, line=1.0)
    label(s, M + 0.44, py + 1.02, 4.4, 0.35, "sinh viên tự tìm đến ChatGPT khi gặp khó khăn",
          12, True, ONBLUE)
    label(s, M, 5.42, 5.3, 1.0,
          [("Sinh viên không thiếu công cụ. Họ thiếu ", MUTED, False),
           ("một kênh an toàn ngay trong lớp", PRIMARY, True),
           (" để nói rằng mình chưa hiểu.", MUTED, False)], 13, line=1.4)

    bx = 6.7
    bw2 = W - M - bx
    card(s, bx, BODY_Y - 0.02, bw2, 4.98)
    pill(s, bx + 0.42, BODY_Y + 0.26, "Survey snapshot", 11, h=0.32)
    pr, pw2 = pill(s, 0, BODY_Y + 0.26, "n = 36", 11, "red", h=0.32)
    pr.left = Inches(bx + bw2 - 0.42 - pw2)

    stats = [("Ngại hỏi vì sợ câu hỏi quá đơn giản", "50%", 0.50, "red"),
             ("Ngại phát biểu trước lớp", "44,4%", 0.444, "red"),
             ("Không biết mô tả lỗi / vấn đề của mình", "41,7%", 0.417, "orange"),
             ("Chờ trên 10 phút hoặc không được hỗ trợ phù hợp", "55,6%", 0.556, "red")]
    y = BODY_Y + 0.9
    for name, val, pct, col in stats:
        label(s, bx + 0.42, y, bw2 - 1.75, 0.34, name, 12, True, PRIMARY_D, line=1.2)
        label(s, bx + bw2 - 1.32, y - 0.07, 0.9, 0.4, val, 19, True,
              ACCENT if col == "red" else ORANGE, align=PP_ALIGN.RIGHT)
        bar(s, bx + 0.42, y + 0.4, bw2 - 0.84, pct, col, 0.15)
        y += 0.85
    label(s, bx + 0.42, y + 0.02, bw2 - 0.84, 0.8,
          "Vấn đề không phải sinh viên không cần giúp đỡ — mà là khoảng trống giữa lúc họ bắt đầu không hiểu "
          "và lúc giảng viên nhận ra.", 12, False, MUTED, line=1.4)
    footer(s, "Evidence")


# ================================================================ 04 LMS context
def slide4():
    s = new_slide()
    topline(s, "03. Bối cảnh LMS", 4)
    h2(s, M, BODY_Y, 11.6, "LMS hiện tại gần như chỉ là nơi mở và đọc PDF.", 30, 0.62)
    sub(s, M, BODY_Y + 0.72, 10.4,
        "Giảng viên thậm chí không dùng nền tảng khi dạy — họ mở file trực tiếp từ máy tính và trình chiếu như bình thường.",
        13.5, 0.55)

    y = 2.95
    cw = (CW - 0.5) / 3
    items = [("PDF", "Chỉ là kho tài liệu",
              "Sinh viên vào LMS để tải và đọc file, không có tương tác nào diễn ra."),
             ("×", "Giảng viên không dùng khi dạy",
              "Quy trình giảng dạy thật đang diễn ra hoàn toàn bên ngoài LMS."),
             ("↺", "Không có vòng phản hồi",
              "Không ai biết lớp đang hiểu tới đâu tại chính thời điểm này.")]
    for i, (g, t, d) in enumerate(items):
        feature(s, M + i * (cw + 0.25), y, cw, 2.15, g, t, d, 15, 12.5)

    card(s, M, 5.4, CW, 1.16, "red")
    label(s, M + 0.5, 5.62, CW - 1.0, 0.8,
          "Vì vậy: không thay toàn bộ LMS, không bắt giảng viên đổi quy trình — chỉ tích hợp thêm một lớp tương tác "
          "vào đúng màn hình học liệu hiện có.", 14.5, True, WHITE, line=1.32)
    footer(s, "Positioning")


# ================================================================ 05 solution
def slide5():
    s = new_slide()
    topline(s, "04. Giải pháp", 5)
    h2(s, M, BODY_Y + 0.1, 5.8, "VLearn\nInteraction Layer", 38, 1.7)
    sub(s, M, BODY_Y + 1.95, 5.5,
        "Một lớp tương tác hai chiều giữa sinh viên và giảng viên, hoạt động song song với file PDF đang trình chiếu.",
        14, 1.0)
    btn(s, M, BODY_Y + 3.15, "Hoạt động ngay trên học liệu hiện có", 13)

    bx = 6.6
    bw2 = W - M - bx
    ch = 2.62
    card(s, bx, BODY_Y - 0.02, bw2, ch, "ice")
    label(s, bx + 0.42, BODY_Y + 0.24, bw2 - 0.84, 0.28, "PHÍA SINH VIÊN", 10.5, True, PRIMARY,
          space=0.8)
    items = ["Phản hồi mức độ hiểu bằng một thao tác",
             "Tham gia quiz hoặc poll nhanh",
             "Xem nội dung cả lớp đang thấy khó",
             "Đặt câu hỏi riêng cho trợ giảng",
             "Bấm “Tôi cũng chưa hiểu” với câu hỏi đã có"]
    y = BODY_Y + 0.66
    for it in items:
        _solid(oval(s, bx + 0.44, y + 0.07, 0.12, 0.12), PRIMARY_L)
        label(s, bx + 0.7, y, bw2 - 1.15, 0.3, it, 12, True, PRIMARY_D)
        y += 0.37

    y2 = BODY_Y + ch + 0.3
    bh = 2.05
    card(s, bx, y2, bw2, bh, "blue")
    label(s, bx + 0.42, y2 + 0.26, bw2 - 0.84, 0.28, "PHÍA GIẢNG VIÊN", 10.5, True, ONBLUE_D,
          space=0.8)
    label(s, bx + 0.42, y2 + 0.62, bw2 - 0.84, 0.4, "Interaction Console", 20, True, WHITE)
    label(s, bx + 0.42, y2 + 1.1, bw2 - 0.84, 0.85,
          "Vẫn mở PDF trực tiếp từ máy như trước. Bên cạnh bài giảng là bảng điều khiển cho thấy trạng thái thực của lớp.",
          12, False, ONBLUE, line=1.35)
    footer(s, "Solution")


# ================================================================ 06 demo pulse
def slide6():
    s = new_slide()
    topline(s, "05. Demo · bước 1", 6)
    h2(s, M, BODY_Y, 5.9, "Learning Pulse:\nmột thao tác,\ncả lớp lên tiếng.", 30, 1.65)
    sub(s, M, BODY_Y + 1.78, 5.75,
        "Sau khi giảng một phần, giảng viên bấm “Gửi Learning Pulse”. Sinh viên nhận ba lựa chọn ngay trong panel "
        "bên cạnh PDF.", 13.5, 0.9)

    y = 3.9
    cw = (5.75 - 0.36) / 3
    for i, (t, k) in enumerate([("Đã hiểu", "blue"), ("Hơi mơ hồ", "orange"), ("Đang bị kẹt", "red")]):
        c = card(s, M + i * (cw + 0.18), y, cw, 0.82, k)
        center_text(c, t, 14, WHITE)

    label(s, M, 5.1, 5.75, 1.2,
          [("Thay vì tiếp tục giảng theo cảm tính, hệ thống đề xuất giảng viên ", MUTED, False),
           ("dừng lại và gửi một poll kiểm tra nhanh", PRIMARY, True),
           (" — sinh viên trả lời ngay trong panel.", MUTED, False)], 13, line=1.42)

    bx = 7.3
    bw2 = W - M - bx
    card(s, bx, BODY_Y - 0.02, bw2, 4.98)
    label(s, bx + 0.42, BODY_Y + 0.26, 3.0, 0.32, "Kết quả Pulse", 16, True, PRIMARY_D)
    pr, pw2 = pill(s, 0, BODY_Y + 0.26, "Live · 36 phản hồi", 10.5, "red", h=0.32)
    pr.left = Inches(bx + bw2 - 0.42 - pw2)

    label(s, bx + 0.42, BODY_Y + 0.8, 2.0, 1.0, "42%", 50, True, ACCENT, line=1.0)
    label(s, bx + 2.15, BODY_Y + 1.0, bw2 - 2.6, 0.75,
          "lớp đang mơ hồ hoặc bị kẹt tại slide hiện tại", 12.5, True, PRIMARY_D, line=1.3)

    rows = [("Đã hiểu", "58%", 0.58, "blue"), ("Hơi mơ hồ", "27%", 0.27, "orange"),
            ("Đang bị kẹt", "15%", 0.15, "red")]
    y2 = BODY_Y + 1.95
    for name, val, pct, col in rows:
        label(s, bx + 0.42, y2, 2.0, 0.28, name, 11.5, True, PRIMARY_D)
        label(s, bx + bw2 - 1.25, y2, 0.83, 0.28, val, 11.5, True, MUTED, align=PP_ALIGN.RIGHT)
        bar(s, bx + 0.42, y2 + 0.3, bw2 - 0.84, pct, col, 0.15)
        y2 += 0.6

    card(s, bx + 0.42, y2 + 0.05, bw2 - 0.84, 1.12, "blue", 0.06)
    label(s, bx + 0.7, y2 + 0.25, bw2 - 1.4, 0.28, "HỆ THỐNG ĐỀ XUẤT", 10, True, ONBLUE_D, space=0.8)
    label(s, bx + 0.7, y2 + 0.57, bw2 - 1.4, 0.55,
          "Dừng lại 3 phút và gửi poll kiểm tra khái niệm vừa giảng.", 12, True, WHITE, line=1.3)
    footer(s, "Learning Pulse")


# ================================================================ 07 demo question
def slide7():
    s = new_slide()
    topline(s, "06. Demo · bước 2", 7)
    h2(s, M, BODY_Y, 11.6, "Hỏi riêng trợ giảng — không cần giải thích lại bối cảnh.", 29, 0.6)
    sub(s, M, BODY_Y + 0.7, 10.8,
        "Một sinh viên chưa hiểu Attention Mask nhưng ngại phát biểu. VLearn tự động đính kèm tên file, trang đang xem "
        "và nội dung liên quan.", 13.5, 0.7)

    y = 2.85
    cw = (CW - 3 * 0.4) / 4
    steps = [("1", "Sinh viên bấm “Hỏi riêng”",
              "Câu hỏi gửi đi kèm file, trang và nội dung đang xem."),
             ("2", "AI gom câu hỏi tương tự",
              "“Có 8 sinh viên chưa hiểu Attention Mask dùng để làm gì.”"),
             ("3", "Giảng viên chọn cách trả lời",
              "Trả lời riêng · trả lời ẩn danh cho cả lớp · đăng Clarification Card."),
             ("4", "Sinh viên xác nhận đã hiểu",
              "Phản hồi đúng ngữ cảnh, ngay tại slide đó.")]
    for i, (n, t, d) in enumerate(steps):
        x = M + i * (cw + 0.4)
        card(s, x, y, cw, 2.35)
        numchip(s, x + 0.26, y + 0.26, n)
        flow(s, x + 0.26, y + 0.82, cw - 0.52, 1.4,
             [(t, 13.5, True, PRIMARY_D, 1.18, 0), (d, 11.5, False, MUTED, 1.32, 7)])
        if i < 3:
            arrow(s, x + cw + 0.07, y + 1.05, 0.26, RGBColor(0xB6, 0xCF, 0xEC))

    card(s, M, 5.5, CW, 1.06, "ice")
    label(s, M + 0.44, 5.74, CW - 0.88, 0.7,
          [("Kết quả: ", PRIMARY, True),
           ("câu hỏi ngại nói ra trở thành dữ liệu giảng viên nhìn thấy được — và trở thành câu trả lời cho cả lớp.",
            BODY, False)], 14, line=1.35)
    footer(s, "Private question → shared clarity")


# ================================================================ 08 console
def slide8():
    s = new_slide()
    topline(s, "07. Interaction Console", 8)
    h2(s, M, BODY_Y - 0.06, 11.6, "Giảng viên nhìn thấy trạng thái thật của lớp.", 29, 0.58)

    bx, by = M, 2.1
    bw2, bh = 7.55, 4.42
    card(s, bx, by, bw2, bh)
    label(s, bx + 0.42, by + 0.3, 4.0, 0.32, "Interaction Console", 16, True, PRIMARY_D)
    pr, pw2 = pill(s, 0, by + 0.3, "Live", 10.5, "red", h=0.32)
    pr.left = Inches(bx + bw2 - 0.42 - pw2)

    metrics = [("36", "sinh viên", TEXT), ("42%", "mơ hồ / kẹt", ACCENT), ("8", "câu hỏi mới", PRIMARY)]
    mw = (bw2 - 0.84 - 0.4) / 3
    for i, (v, t, c) in enumerate(metrics):
        x = bx + 0.42 + i * (mw + 0.2)
        card(s, x, by + 0.8, mw, 0.95, "ice", 0.05)
        label(s, x + 0.22, by + 0.9, mw - 0.44, 0.42, v, 22, True, c)
        label(s, x + 0.22, by + 1.36, mw - 0.44, 0.28, t, 10.5, True, MUTED)

    label(s, bx + 0.42, by + 2.0, 3.2, 0.28, "PULSE THEO SLIDE", 10, True, SMALLC, space=0.8)
    hx, hy = bx + 0.42, by + 2.35
    cell, gap = 0.235, 0.07
    pattern = [1, 1, 2, 1, 3, 3, 1, 2, 1, 1, 3, 2,
               1, 2, 3, 3, 1, 1, 2, 1, 3, 1, 1, 2,
               1, 1, 1, 2, 3, 3, 3, 1, 2, 1, 1, 1]
    colmap = {1: RGBColor(0x17, 0x79, 0xDF), 2: ORANGE, 3: ACCENT}
    for i, v in enumerate(pattern):
        r, c = divmod(i, 12)
        _solid(rect(s, hx + c * (cell + gap), hy + r * (cell + gap), cell, cell, 0.28), colmap[v])
    lx = hx
    for i, t in enumerate(["Đã hiểu", "Hơi mơ hồ", "Đang bị kẹt"]):
        _solid(oval(s, lx, by + 3.46, 0.13, 0.13), colmap[i + 1])
        label(s, lx + 0.2, by + 3.42, 1.4, 0.26, t, 10.5, True, SMALLC)
        lx += 1.3

    label(s, bx + 4.35, by + 2.0, 3.0, 0.28, "CỤM CÂU HỎI (AI)", 10, True, SMALLC, space=0.8)
    qs = [("Attention Mask dùng để làm gì?", "8 SV", ACCENT),
          ("Khác nhau giữa encoder / decoder", "5 SV", ORANGE),
          ("Vì sao cần positional encoding", "3 SV", PRIMARY)]
    qy = by + 2.3
    for t, n, col in qs:
        card(s, bx + 4.35, qy, 2.78, 0.6, "ice", 0.05)
        label(s, bx + 4.52, qy + 0.08, 1.85, 0.46, t, 10.5, True, PRIMARY_D, line=1.2)
        label(s, bx + 6.35, qy + 0.16, 0.6, 0.28, n, 11, True, col, align=PP_ALIGN.RIGHT)
        qy += 0.7

    rx = 8.62
    rw = W - M - rx
    card(s, rx, by, rw, 1.75, "red")
    label(s, rx + 0.36, by + 0.26, rw - 0.72, 0.28, "CẢNH BÁO", 10, True, ONRED, space=0.8)
    label(s, rx + 0.36, by + 0.62, rw - 0.72, 1.0,
          "42% lớp chưa theo kịp tại slide 14 — nên dừng lại trước khi sang phần mới.",
          13.5, True, WHITE, line=1.3)

    card(s, rx, by + 1.97, rw, 2.45, "blue")
    label(s, rx + 0.36, by + 2.22, rw - 0.72, 0.28, "CLARIFICATION CARD", 10, True, ONBLUE_D,
          space=0.8)
    label(s, rx + 0.36, by + 2.56, rw - 0.72, 0.34, "Attention Mask", 17, True, WHITE)
    label(s, rx + 0.36, by + 2.98, rw - 0.72, 0.9,
          "“Mask dùng để chặn mô hình nhìn vào các token không được phép — ví dụ token tương lai.”",
          11.5, False, ONBLUE, line=1.35)
    btn(s, rx + 0.36, by + 3.78, "Đăng ngay tại slide 14", 11.5, 0.42)
    footer(s, "Teacher view")


# ================================================================ 09 AI role
def slide9():
    s = new_slide()
    topline(s, "08. Vai trò của AI", 9)
    h2(s, M, BODY_Y, 6.0, "AI đứng phía sau,\nkhông đứng phía trước.", 30, 1.15)
    sub(s, M, BODY_Y + 1.3, 5.8,
        "Chúng em không thêm một chatbot nữa vào lớp học. AI xử lý tín hiệu để giảng viên ra quyết định nhanh hơn.",
        13.5, 0.9)

    y = 3.3
    items = [("∑", "Gom các câu hỏi giống nhau", "Nhiều câu hỏi rời rạc trở thành một vấn đề chung."),
             ("!", "Phát hiện nội dung gây khó", "Biết chính xác slide nào đang làm lớp chững lại."),
             ("▦", "Tóm tắt tín hiệu từ cả lớp", "Pulse, poll, câu hỏi được gói lại thành một bức tranh."),
             ("↗", "Đề xuất thời điểm can thiệp", "Gợi ý khi nào nên dừng lại, khi nào nên đi tiếp.")]
    cwid = 2.85
    for i, (g, t, d) in enumerate(items):
        r, c = divmod(i, 2)
        x = M + c * (cwid + 0.2)
        yy = y + r * 1.42
        card(s, x, yy, cwid, 1.38)
        icon(s, x + 0.22, yy + 0.22, g, 0.4, fsize=14)
        flow(s, x + 0.75, yy + 0.22, cwid - 0.97, 1.04,
             [(t, 12.5, True, PRIMARY_D, 1.15, 0), (d, 10.5, False, MUTED, 1.3, 5)])

    bx = 7.05
    bw2 = W - M - bx
    card(s, bx, BODY_Y, bw2, 3.05, "red")
    label(s, bx + 0.5, BODY_Y + 0.34, bw2 - 1.0, 0.28, "GIỚI HẠN CỦA AI", 10.5, True, ONRED,
          space=0.8)
    yy = BODY_Y + 0.82
    for t in ["Không trả lời thay giảng viên", "Không tự chấm điểm cuối cùng",
              "Không âm thầm đánh giá sinh viên"]:
        label(s, bx + 0.5, yy - 0.06, 0.4, 0.42, "×", 20, True, WHITE)
        label(s, bx + 0.92, yy, bw2 - 1.42, 0.36, t, 14.5, True, WHITE)
        yy += 0.62

    card(s, bx, BODY_Y + 3.35, bw2, 1.55, "ice")
    label(s, bx + 0.5, BODY_Y + 3.62, bw2 - 1.0, 1.0,
          [("Quyết định cuối cùng ", BODY, False), ("vẫn thuộc về giảng viên.", PRIMARY, True)],
          18, line=1.3)
    footer(s, "Human-in-the-loop")


# ================================================================ 10 priority
def slide10():
    s = new_slide()
    topline(s, "09. Ưu tiên tính năng", 10)
    h2(s, M, BODY_Y, 5.9, "MVP bám theo thứ tự\nưu tiên của người học.", 29, 1.15)
    sub(s, M, BODY_Y + 1.28, 5.75,
        "Khảo sát cho thấy nhu cầu tập trung vào hỗ trợ đúng lúc, không phải vào phần thưởng hay điểm thi đua.",
        13.5, 0.85)

    y = 3.3
    for i, t in enumerate(["Không bảng xếp hạng", "Không badge", "Không AI làm mọi thứ"]):
        yy = y + i * 0.56
        o = rect(s, M, yy, 0.34, 0.34, 0.3)
        _solid(o, RGBColor(0xFF, 0xE4, 0xE5))
        center_text(o, "×", 14, ACCENT)
        label(s, M + 0.5, yy + 0.05, 3.6, 0.3, t, 13, True, MUTED)

    label(s, M, y + 1.8, 5.75, 0.8,
          "MVP chỉ tập trung vào vòng lặp quan trọng nhất của một tiết học.", 13, True, PRIMARY,
          line=1.35)

    bx = 7.05
    bw2 = W - M - bx
    card(s, bx, BODY_Y - 0.02, bw2, 3.98)
    label(s, bx + 0.42, BODY_Y + 0.26, bw2 - 0.84, 0.28, "MỨC ĐỘ ƯU TIÊN THEO KHẢO SÁT", 10,
          True, SMALLC, space=0.8)
    rows = [("Gợi ý / hỗ trợ cá nhân khi gặp lỗi", 0.95, "blue", "Cao nhất"),
            ("Cảnh báo giảng viên khi lớp gặp khó", 0.82, "blue", "Cao"),
            ("Quiz & poll nhanh trong giờ học", 0.66, "orange", "Rõ ràng"),
            ("Gamification (badge, xếp hạng)", 0.18, "grey", "Rất thấp")]
    yy = BODY_Y + 0.72
    for name, pct, col, tag in rows:
        label(s, bx + 0.42, yy, bw2 - 1.65, 0.3, name, 12, True,
              MUTED if col == "grey" else PRIMARY_D)
        label(s, bx + bw2 - 1.25, yy + 0.02, 0.83, 0.28, tag, 10.5, True,
              MUTED if col == "grey" else PRIMARY, align=PP_ALIGN.RIGHT)
        bar(s, bx + 0.42, yy + 0.36, bw2 - 0.84, pct, col, 0.16)
        yy += 0.85

    card(s, M, 5.5, CW, 1.0, "blue")
    steps = ["Sinh viên thể hiện khó khăn", "Giảng viên nhìn thấy", "Tương tác xảy ra",
             "Sinh viên xác nhận đã hiểu"]
    sw = (CW - 0.8) / 4
    for i, t in enumerate(steps):
        x = M + 0.4 + i * sw
        label(s, x, 5.78, sw - 0.5, 0.5, t, 12, True, WHITE, align=PP_ALIGN.CENTER, line=1.2)
        if i < 3:
            arrow(s, x + sw - 0.42, 5.9, 0.24, RGBColor(0x8C, 0xC2, 0xFF))
    footer(s, "MVP scope")


# ================================================================ 11 impact
def slide11():
    s = new_slide()
    topline(s, "10. Giá trị", 11)
    h2(s, M, BODY_Y, 11.6, "Một màn hình đọc PDF thụ động trở thành\nkhông gian học có tương tác.", 29, 1.05)
    sub(s, M, BODY_Y + 1.22, 10.6,
        "Không thay đổi thói quen giảng dạy hiện tại, không yêu cầu nhà trường xây lại hệ thống.", 13.5, 0.5)

    y = 2.98
    cw = (CW - 0.5) / 3
    cards = [("Sinh viên", "Một cách an toàn để nói rằng mình chưa hiểu — không cần giơ tay trước cả lớp.", "white"),
             ("Giảng viên", "Nhìn thấy trạng thái thật của lớp trước khi vấn đề biến thành điểm thấp.", "blue"),
             ("Nhà trường", "Nâng cấp LMS hiện có mà không cần xây lại toàn bộ hệ thống.", "white")]
    for i, (t, d, k) in enumerate(cards):
        x = M + i * (cw + 0.25)
        card(s, x, y, cw, 2.14, k)
        if k == "blue":
            chip = rect(s, x + 0.34, y + 0.3, 0.44, 0.44, 0.3)
            _solid(chip, PRIMARY_L)
            center_text(chip, "✓", 14, WHITE)
            tcol, dcol = WHITE, ONBLUE
        else:
            icon(s, x + 0.34, y + 0.3, "✓", 0.44, fsize=14)
            tcol, dcol = PRIMARY_D, MUTED
        flow(s, x + 0.34, y + 0.88, cw - 0.68, 1.3,
             [(t, 18, True, tcol, 1.15, 0), (d, 12, False, dcol, 1.35, 8)])

    card(s, M, 5.5, CW, 1.08, "orange")
    label(s, M + 0.5, 5.74, 3.0, 0.4, "North Star", 18, True, WHITE)
    label(s, M + 3.7, 5.76, CW - 4.2, 0.6,
          "Tỷ lệ sinh viên xác nhận “đã hiểu” ngay trong buổi học, thay vì mang thắc mắc về nhà.",
          14, True, WHITE, line=1.3)
    footer(s, "Impact")


# ================================================================ 12 message
def slide12():
    s = new_slide(wave=True)
    topline(s, "11. Thông điệp", 12)
    dots(s, W - 2.5, 1.15, 7, 4)

    label(s, M, 2.05, 10.6, 0.95, "“", 56, True, RGBColor(0xC9, 0xDF, 0xF9), line=1.0)
    flow(s, M, 2.72, 11.6, 1.5,
         [("Sinh viên không cần thêm một nơi để đọc tài liệu.", 28, True, TEXT, 1.2, 0),
          ([("Họ cần một cách để ", TEXT, True),
            ("được lắng nghe ngay trong lúc học.", PRIMARY, True)], 28, True, TEXT, 1.2, 8)])

    y = 4.7
    x = M
    pills = ["Learning Pulse", "Poll & Quiz nhanh", "Hỏi riêng trợ giảng", "AI gom câu hỏi",
             "Clarification Card"]
    for i, t in enumerate(pills):
        p, pw = pill(s, x, y, t, 12, "light" if i % 2 == 0 else "ghost", h=0.42)
        x += pw + 0.18

    label(s, M, 5.55, 10.9, 0.5,
          "VLearn Interaction Layer — lớp tương tác đặt ngay trên học liệu hiện có.", 14, True, MUTED)
    footer(s, "Thank you")


# ================================================================ 13 summary
def slide13():
    s = new_slide()
    topline(s, "12. Tóm tắt 5 phút", 13)
    h2(s, M, BODY_Y - 0.05, 11.6, "Toàn bộ câu chuyện trong một trang.", 29, 0.58)

    y = 2.2
    cw = (CW - 3 * 0.25) / 4
    blocks = [("VẤN ĐỀ", "Sinh viên không nói ra khi chưa hiểu; giảng viên nhận ra quá muộn.", "white"),
              ("BẰNG CHỨNG", "Khảo sát 36 SV: 88,9% tự hỏi ChatGPT; 55,6% chờ hơn 10 phút hoặc không được hỗ trợ.", "white"),
              ("GIẢI PHÁP", "Lớp tương tác hai chiều cạnh PDF: pulse, poll, hỏi riêng, clarification card.", "blue"),
              ("VÌ SAO KHẢ THI", "Không thay LMS, không đổi thói quen dạy; AI chỉ hỗ trợ phía sau.", "white")]
    for i, (t, d, k) in enumerate(blocks):
        x = M + i * (cw + 0.25)
        card(s, x, y, cw, 2.15, k)
        tcol = ONBLUE_D if k == "blue" else PRIMARY
        dcol = WHITE if k == "blue" else BODY
        flow(s, x + 0.3, y + 0.32, cw - 0.6, 1.6,
             [(t, 10.5, True, tcol, 1.2, 0, PP_ALIGN.LEFT, 0.8),
              (d, 12.5, False, dcol, 1.38, 10)])

    card(s, M, 4.72, CW, 1.7, "ice")
    label(s, M + 0.5, 4.98, CW - 1.0, 0.28, "VÒNG LẶP LÕI", 10.5, True, SMALLC, space=0.8)
    steps = ["Sinh viên thể hiện khó khăn", "Giảng viên nhìn thấy", "Tương tác xảy ra",
             "Sinh viên xác nhận đã hiểu"]
    sw = (CW - 1.0) / 4
    for i, t in enumerate(steps):
        x = M + 0.5 + i * sw
        label(s, x, 5.44, sw - 0.42, 0.7, t, 13.5, True, PRIMARY_D, line=1.25)
        if i < 3:
            arrow(s, x + sw - 0.46, 5.56, 0.28, PRIMARY)
    footer(s, "Summary")


for fn in (slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8, slide9,
           slide10, slide11, slide12, slide13):
    fn()

# ---------------------------------------------------------------- speaker notes
SCRIPT = [
    """[0:00–0:10] Mở đầu
Em xin bắt đầu bằng một tình huống rất quen thuộc. Đây là VLearn Interaction Layer — lớp tương tác
đặt ngay trên màn hình học liệu hiện có.""",
    """[0:10–0:35] Vấn đề
Trong một lớp học đông, giảng viên đang trình chiếu PDF và tiếp tục giảng. Ở phía dưới, có những sinh viên
đã hiểu, có người chỉ hiểu một phần, và cũng có người đã bị kẹt từ vài slide trước.
Nhưng gần như không ai nói ra. Sinh viên ngại phát biểu, sợ câu hỏi của mình quá đơn giản, hoặc thậm chí
không biết phải mô tả điều mình chưa hiểu như thế nào. Còn giảng viên nhìn cả lớp im lặng và mặc định
rằng mọi người vẫn đang theo kịp.""",
    """[0:35–1:15] Bằng chứng
Chúng em đã khảo sát 36 sinh viên có trải nghiệm học thực hành trên LMS. Kết quả cho thấy 88,9% sinh viên
thường tự tìm đến ChatGPT khi gặp khó khăn. Tuy nhiên, 50% ngại hỏi vì sợ câu hỏi quá đơn giản,
44,4% ngại phát biểu trước lớp và 41,7% không biết mô tả lỗi hoặc vấn đề của mình. Khoảng 55,6% sinh viên
phải chờ trên 10 phút hoặc thường không nhận được hỗ trợ phù hợp.
Như vậy, vấn đề không phải sinh viên không cần giúp đỡ. Vấn đề là giữa lúc sinh viên bắt đầu không hiểu
và lúc giảng viên nhận ra điều đó đang tồn tại một khoảng trống rất lớn.""",
    """[1:15–1:45] Bối cảnh LMS hiện tại
LMS hiện tại gần như chỉ là nơi sinh viên mở và đọc PDF. Giảng viên thậm chí không sử dụng nền tảng này
khi dạy — họ vẫn mở file trực tiếp từ máy tính và trình chiếu như bình thường.
Vì vậy, chúng em không cố thay đổi toàn bộ LMS, cũng không yêu cầu giảng viên chuyển sang một quy trình
giảng dạy mới. Chúng em chỉ tích hợp thêm một lớp tương tác vào đúng màn hình học liệu hiện có.""",
    """[1:45–2:20] Giới thiệu giải pháp
Giải pháp của chúng em là VLearn Interaction Layer — một lớp tương tác hai chiều giữa sinh viên và
giảng viên, hoạt động song song với file PDF.
Ở phía sinh viên: phản hồi mức độ hiểu bằng một thao tác; tham gia quiz hoặc poll nhanh; xem những nội dung
cả lớp đang cảm thấy khó; đặt câu hỏi riêng cho trợ giảng; hoặc bấm "Tôi cũng chưa hiểu" với một câu hỏi đã có.
Ở phía giảng viên, họ vẫn mở PDF trực tiếp từ máy như trước. Nhưng bên cạnh bài giảng sẽ có một
Interaction Console giúp họ nhìn thấy trạng thái thực của lớp.""",
    """[2:20–2:50] Demo — bước 1
Giảng viên mở file PDF trên máy và bắt đầu phiên học. Sau khi giảng một phần, giảng viên bấm
"Gửi Learning Pulse". Ngay lập tức, sinh viên nhận ba lựa chọn: "Đã hiểu", "Hơi mơ hồ", hoặc "Đang bị kẹt".
Giả sử kết quả trả về cho thấy 42% lớp đang mơ hồ hoặc bị kẹt. Thay vì tiếp tục giảng theo cảm tính,
hệ thống đề xuất giảng viên dừng lại và gửi một poll kiểm tra nhanh. Sinh viên trả lời poll ngay trong
panel bên cạnh PDF.""",
    """[2:50–3:20] Demo — bước 2
Đồng thời, một sinh viên chưa hiểu Attention Mask nhưng ngại phát biểu. Bạn ấy bấm "Hỏi riêng trợ giảng".
VLearn tự động đính kèm tên file, trang đang xem và nội dung liên quan — sinh viên không cần phải giải thích
lại toàn bộ bối cảnh.
Nếu có nhiều câu hỏi tương tự, AI chạy phía sau để gom chúng thành một vấn đề chung, ví dụ:
"Có 8 sinh viên chưa hiểu Attention Mask dùng để làm gì."
Giảng viên có thể trả lời riêng, trả lời ẩn danh cho cả lớp, hoặc đăng một Clarification Card ngay tại slide đó.""",
    """[3:20–3:30] Demo — màn hình giảng viên
Đây là Interaction Console. Sinh viên nhận phản hồi đúng ngữ cảnh, xác nhận "Đã hiểu", và giảng viên mới
tiếp tục sang phần tiếp theo.""",
    """[3:30–4:05] Vai trò của AI
Điểm quan trọng là chúng em không đặt AI ở phía trước dưới dạng một chatbot khác. AI hoạt động ở phía sau để:
gom các câu hỏi giống nhau; phát hiện nội dung đang gây khó khăn; tóm tắt tín hiệu từ cả lớp; và đề xuất
thời điểm giảng viên nên can thiệp.
Quyết định cuối cùng vẫn thuộc về giảng viên. AI không thay giảng viên trả lời, không tự chấm điểm cuối cùng
và không âm thầm đánh giá sinh viên.""",
    """[4:05–4:35] Vì sao ưu tiên các tính năng này?
Trong khảo sát, tính năng được chọn quan trọng nhất là gợi ý hoặc hỗ trợ cá nhân khi gặp lỗi, tiếp theo là
cảnh báo giảng viên khi cả lớp gặp khó khăn. Quiz và poll nhanh cũng nhận được nhu cầu rõ ràng. Trong khi đó,
gamification có mức ưu tiên rất thấp.
Vì vậy, MVP của chúng em không có bảng xếp hạng, không có badge và không cố xây một trợ lý AI làm mọi thứ.
Chúng em chỉ tập trung vào vòng lặp quan trọng nhất: sinh viên thể hiện khó khăn → giảng viên nhìn thấy →
tương tác xảy ra → sinh viên xác nhận đã hiểu.""",
    """[4:35–4:50] Giá trị
VLearn Interaction Layer biến một màn hình đọc PDF thụ động thành một không gian học có tương tác, nhưng
không làm thay đổi thói quen giảng dạy hiện tại.
Với sinh viên, đây là một cách an toàn để nói rằng mình chưa hiểu. Với giảng viên, đây là khả năng nhìn thấy
trạng thái thật của lớp trước khi vấn đề biến thành điểm thấp hoặc sinh viên bị bỏ lại phía sau. Và với
nhà trường, đây là một cách nâng cấp LMS hiện có mà không cần xây lại toàn bộ hệ thống.""",
    """[4:50–5:00] Kết thúc
Thông điệp của chúng em rất đơn giản: Sinh viên không cần thêm một nơi để đọc tài liệu.
Họ cần một cách để được lắng nghe ngay trong lúc học.
Em xin cảm ơn.""",
    """[Slide dự phòng — Q&A]
Dùng khi ban giám khảo hỏi lại tổng quan: vấn đề, bằng chứng, giải pháp, tính khả thi và vòng lặp lõi
của sản phẩm.""",
]
for sld, txt in zip(prs.slides, SCRIPT):
    notes(sld, txt)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("saved:", OUT)
if WARN:
    print("\n".join(WARN))
    print("total warnings:", len(WARN))
else:
    print("no overflow warnings")
